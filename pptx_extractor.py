"""
PPTX extraction using python-pptx.

Extracts slide titles, text frames, tables, images, and speaker notes
from .pptx files, outputting the same PageResult format as pdf_extractor.py
so html_builder can consume them unchanged.

Each slide maps to one PageResult (page_number = slide index, 0-based).
Speaker notes are marked with a special font name so html_builder can
render them as accessible associated content.
"""

import io
import logging

from pptx import Presentation

from pdf_extractor import TextSpan, ImageInfo, TableData, PageResult

logger = logging.getLogger(__name__)

# Default slide dimensions (standard 4:3 in points: 10" × 7.5")
_DEFAULT_WIDTH = 720.0
_DEFAULT_HEIGHT = 540.0

# Synthetic layout constants (points)
_LEFT_MARGIN = 72.0
_LINE_HEIGHT = 16.0

# Font sizes that map to html_builder's _heading_level() thresholds:
#   h1 >= 24.0, h2 >= 18.0
_SLIDE_TITLE_SIZE = 26.0   # detected as h1 by html_builder

# Default body text size (points)
_BODY_SIZE = 11.0

# Image content-type to extension map
_IMG_EXT_MAP: dict[str, str] = {
    "image/png": "png",
    "image/jpeg": "jpeg",
    "image/jpg": "jpeg",
    "image/gif": "gif",
    "image/bmp": "bmp",
    "image/tiff": "tiff",
}

# Speaker notes marker — html_builder detects this font name and renders
# the span as an accessible <aside> / <details> element.
SPEAKER_NOTES_FONT = "SpeakerNotes"


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_pptx(file_bytes: bytes) -> tuple[list[PageResult], dict]:
    """Extract content from a .pptx file.

    Returns
    -------
    pages : list[PageResult]
        One PageResult per slide (page_number = slide index, 0-based).
    metadata : dict
        Presentation metadata (title, author, subject, format).
    """
    prs = Presentation(io.BytesIO(file_bytes))

    # --- slide dimensions ------------------------------------------------
    width = _DEFAULT_WIDTH
    height = _DEFAULT_HEIGHT
    if prs.slide_width is not None:
        width = prs.slide_width / 12700.0  # EMU → points
    if prs.slide_height is not None:
        height = prs.slide_height / 12700.0

    # --- metadata --------------------------------------------------------
    metadata: dict = {"format": "pptx"}
    try:
        cp = prs.core_properties
        metadata["title"] = cp.title or ""
        metadata["author"] = cp.author or ""
        metadata["subject"] = cp.subject or ""
    except Exception:
        pass

    # --- extract slides --------------------------------------------------
    pages: list[PageResult] = []

    for slide_idx, slide in enumerate(prs.slides):
        text_spans: list[TextSpan] = []
        images: list[ImageInfo] = []
        tables: list[TableData] = []
        y = 50.0  # synthetic vertical cursor

        # --- slide title -------------------------------------------------
        title_shape = slide.shapes.title
        title_text = ""
        if title_shape is not None:
            title_text = (title_shape.text or "").strip()

        if title_text:
            text_spans.append(TextSpan(
                text=title_text,
                x0=_LEFT_MARGIN,
                y0=y,
                x1=_LEFT_MARGIN + len(title_text) * 10,
                y1=y + _SLIDE_TITLE_SIZE,
                font="Arial",
                size=_SLIDE_TITLE_SIZE,
                color=0,
                bold=True,
                italic=False,
            ))
            y += _SLIDE_TITLE_SIZE + 12.0

        # --- shapes (text, tables, images) --------------------------------
        title_id = title_shape.shape_id if title_shape is not None else None

        # Sort non-title shapes by position for natural reading order
        sorted_shapes = sorted(
            (s for s in slide.shapes
             if title_id is None or s.shape_id != title_id),
            key=lambda s: (int(s.top or 0), int(s.left or 0)),
        )

        for shape in sorted_shapes:
            # Tables
            if shape.has_table:
                tbl = shape.table
                all_rows: list[list[str]] = []
                for row in tbl.rows:
                    all_rows.append([cell.text.strip() for cell in row.cells])
                if all_rows:
                    header = all_rows[0]
                    data = all_rows[1:] if len(all_rows) > 1 else []
                    tables.append(TableData(
                        bbox=(_LEFT_MARGIN, y, 540.0, y + 100.0),
                        header=header,
                        rows=data,
                    ))
                    y += 100.0 + 8.0
                continue

            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        y += _LINE_HEIGHT * 0.5
                        continue

                    # Determine font properties from runs
                    font_size = _BODY_SIZE
                    is_bold = False
                    is_italic = False
                    runs_with_text = [r for r in para.runs if r.text.strip()]
                    if runs_with_text:
                        # Use first explicit font size
                        for run in runs_with_text:
                            if run.font.size is not None:
                                font_size = run.font.size / 12700.0
                                break
                        is_bold = all(bool(r.font.bold) for r in runs_with_text)
                        is_italic = all(bool(r.font.italic) for r in runs_with_text)

                    text_spans.append(TextSpan(
                        text=text,
                        x0=_LEFT_MARGIN,
                        y0=y,
                        x1=_LEFT_MARGIN + len(text) * 6,
                        y1=y + font_size,
                        font="Arial",
                        size=font_size,
                        color=0,
                        bold=is_bold,
                        italic=is_italic,
                    ))
                    y += font_size + 4.0
                continue

            # Images — try to extract from picture shapes
            try:
                img = shape.image
                img_bytes = img.blob
                content_type = img.content_type
                ext = _IMG_EXT_MAP.get(content_type, "png")

                sx0 = float(shape.left or 0) / 12700.0
                sy0 = float(shape.top or 0) / 12700.0
                sw = float(shape.width or 0) / 12700.0
                sh = float(shape.height or 0) / 12700.0

                images.append(ImageInfo(
                    page_number=slide_idx,
                    x0=sx0,
                    y0=sy0,
                    x1=sx0 + sw,
                    y1=sy0 + sh,
                    image_bytes=img_bytes,
                    extension=ext,
                    xref=len(images),
                ))
            except (AttributeError, TypeError):
                pass  # Not a picture shape
            except Exception:
                logger.warning(
                    "Failed to extract image from slide %d", slide_idx + 1
                )

        # --- speaker notes -----------------------------------------------
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_spans.append(TextSpan(
                        text=notes_text,
                        x0=_LEFT_MARGIN,
                        y0=y,
                        x1=_LEFT_MARGIN + len(notes_text) * 6,
                        y1=y + _BODY_SIZE,
                        font=SPEAKER_NOTES_FONT,
                        size=_BODY_SIZE,
                        color=0,
                        bold=False,
                        italic=False,
                    ))
                    y += _BODY_SIZE + 4.0
            except Exception:
                logger.warning(
                    "Failed to extract notes from slide %d", slide_idx + 1
                )

        page = PageResult(
            page_number=slide_idx,
            width=width,
            height=height,
            is_scanned=False,  # PPTX is always digital
            text_spans=text_spans,
            images=images,
            tables=tables,
        )
        pages.append(page)

    return pages, metadata
