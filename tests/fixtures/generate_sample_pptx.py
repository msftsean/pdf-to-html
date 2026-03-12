"""
Generate tests/fixtures/sample.pptx programmatically (T067).

Creates a 5-slide PowerPoint deck with:
- Slide 1: Title slide (title + subtitle)
- Slide 2: Content slide with bullet points and speaker notes
- Slide 3: Content slide with a data table
- Slide 4: Content slide with an embedded image
- Slide 5: Mixed content slide with table, text, and speaker notes

Run:  python tests/fixtures/generate_sample_pptx.py
"""

import io
import os
import struct
import zlib

from pptx import Presentation
from pptx.util import Inches, Pt


def _create_small_png(width: int = 10, height: int = 10) -> bytes:
    """Create a minimal valid RGB PNG (solid red)."""
    raw = b""
    for _ in range(height):
        raw += b"\x00"  # filter byte: None
        for _ in range(width):
            raw += b"\xff\x00\x00"  # red pixel
    compressed = zlib.compress(raw)

    def _chunk(tag: bytes, data: bytes) -> bytes:
        c = tag + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", compressed)
        + _chunk(b"IEND", b"")
    )


def generate_sample_pptx() -> Presentation:
    """Build a sample 5-slide PPTX with all major content types."""
    prs = Presentation()
    prs.core_properties.title = "NCDIT Sample Presentation"
    prs.core_properties.author = "Test Suite"

    # --- Slide 1: Title Slide -------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Department of Information Technology"
    slide1.placeholders[1].text = "Annual Report 2025"

    # --- Slide 2: Content with bullets and speaker notes ----------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Strategic Priorities"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Modernize legacy systems"
    for item in [
        "Improve cybersecurity posture",
        "Expand digital services",
        "Train workforce on new technologies",
    ]:
        p = tf2.add_paragraph()
        p.text = item

    # Add speaker notes
    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = "Discuss each priority and timeline for implementation."

    # --- Slide 3: Content with data table --------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide3.shapes.title.text = "Budget Overview"

    table_shape = slide3.shapes.add_table(
        rows=4, cols=3,
        left=Inches(1.0), top=Inches(2.0),
        width=Inches(8.0), height=Inches(3.0),
    )
    tbl = table_shape.table
    for row_idx, (cat, amt, pct) in enumerate([
        ("Category", "Amount", "Percentage"),
        ("Personnel", "$1.5M", "65%"),
        ("Operations", "$0.5M", "22%"),
        ("Technology", "$0.3M", "13%"),
    ]):
        tbl.cell(row_idx, 0).text = cat
        tbl.cell(row_idx, 1).text = amt
        tbl.cell(row_idx, 2).text = pct

    # --- Slide 4: Content with embedded image ----------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide4.shapes.title.text = "Infrastructure Diagram"

    png_bytes = _create_small_png()
    slide4.shapes.add_picture(
        io.BytesIO(png_bytes),
        left=Inches(2.0), top=Inches(2.0),
        width=Inches(4.0), height=Inches(3.0),
    )

    # --- Slide 5: Mixed content with speaker notes -----------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Review the budget allocation for Q3."
    p5 = tf5.add_paragraph()
    p5.text = "Schedule follow-up meetings with department heads."

    notes5 = slide5.notes_slide.notes_text_frame
    notes5.text = "Emphasize the importance of timely budget submissions."

    return prs


if __name__ == "__main__":
    fixtures_dir = os.path.dirname(os.path.abspath(__file__))
    os.makedirs(fixtures_dir, exist_ok=True)
    out = os.path.join(fixtures_dir, "sample.pptx")
    generate_sample_pptx().save(out)
    print(f"Generated: {out}")
