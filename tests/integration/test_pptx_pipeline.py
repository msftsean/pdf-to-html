"""
Integration test: PPTX → extract → html_builder → WCAG validation (T066).

Verifies the full pipeline from a PowerPoint presentation through extraction
and HTML building to WCAG compliance checking.  Confirms that the
html_builder produces accessible output from PPTX-sourced PageResult
objects with zero critical/serious violations.
"""

import io
import re
import struct
import zlib

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from backend.pptx_extractor import extract_pptx, SPEAKER_NOTES_FONT
from backend.html_builder import build_html
from backend.wcag_validator import validate_html


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _pptx_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _create_small_png(width: int = 4, height: int = 4) -> bytes:
    """Create a minimal valid RGB PNG (solid red)."""
    raw = b""
    for _ in range(height):
        raw += b"\x00"
        for _ in range(width):
            raw += b"\xff\x00\x00"
    compressed = zlib.compress(raw)

    def _chunk(tag: bytes, data: bytes) -> bytes:
        c = tag + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", compressed)
        + _chunk(b"IEND", b"")
    )


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

class TestPptxPipeline:
    """Full PPTX pipeline: extract → build_html → wcag_validator."""

    def test_simple_pptx_zero_critical_serious(self):
        """Basic PPTX with title slide should pass WCAG."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Annual Report"

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)
        violations = validate_html(html_out)

        critical_serious = [
            v for v in violations if v.severity in ("critical", "serious")
        ]
        assert critical_serious == [], (
            f"WCAG violations: "
            f"{[(v.rule_id, v.description) for v in critical_serious]}"
        )

    def test_html_has_lang(self):
        """Output HTML must have lang attribute."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)
        assert 'lang="en"' in html_out

    def test_html_has_skip_nav_and_main(self):
        """Output must contain skip-nav link and main landmark."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)
        assert 'class="skip-nav"' in html_out
        assert 'id="main-content"' in html_out

    def test_slide_sections_have_aria_label(self):
        """Each slide renders as a <section> with aria-label='Slide N'."""
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "First"
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2.shapes.title.text = "Second"

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)

        assert 'aria-label="Slide 1: First"' in html_out
        assert 'aria-label="Slide 2: Second"' in html_out

    def test_slide_title_rendered_as_heading(self):
        """Slide title appears as <h1> in generated HTML."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Department Overview"

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)

        assert "<h1>" in html_out
        assert "Department Overview" in html_out

    def test_table_accessible(self):
        """Extracted tables produce <th scope='col'> in HTML."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Data"
        tbl_shape = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(1), top=Inches(2),
            width=Inches(6), height=Inches(2),
        )
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Name"
        tbl.cell(0, 1).text = "Value"
        tbl.cell(1, 0).text = "X"
        tbl.cell(1, 1).text = "42"

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)

        assert "<table" in html_out
        assert '<th scope="col">' in html_out
        assert "Name" in html_out
        assert "42" in html_out

        # No table-related WCAG violations
        violations = validate_html(html_out)
        table_violations = [
            v for v in violations
            if "table" in v.rule_id.lower() or "th" in v.rule_id.lower()
        ]
        assert table_violations == []

    def test_image_accessible(self):
        """Embedded images have alt text and figure/figcaption."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Photos"
        slide.shapes.add_picture(
            io.BytesIO(_create_small_png()),
            left=Inches(1), top=Inches(2),
            width=Inches(2), height=Inches(2),
        )

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta, embed_images=True)

        assert "<figure>" in html_out
        assert "alt=" in html_out

        violations = validate_html(html_out)
        img_violations = [v for v in violations if v.rule_id == "image-alt"]
        assert img_violations == []

    def test_speaker_notes_rendered_as_aside(self):
        """Speaker notes render as <aside> with <details> for accessibility."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Important Slide"
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Discuss the key findings."

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)

        assert '<aside class="speaker-notes"' in html_out
        assert "Speaker Notes" in html_out
        assert "Discuss the key findings" in html_out
        assert "<details>" in html_out
        assert "role=\"note\"" in html_out

    def test_speaker_notes_wcag_compliant(self):
        """Slides with speaker notes still pass WCAG validation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title Slide"
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Speaker notes content."

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)
        violations = validate_html(html_out)

        critical_serious = [
            v for v in violations if v.severity in ("critical", "serious")
        ]
        assert critical_serious == [], (
            f"WCAG violations with speaker notes: "
            f"{[(v.rule_id, v.description) for v in critical_serious]}"
        )

    def test_no_speaker_notes_no_aside(self):
        """Slides without speaker notes don't render <aside>."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "No Notes"

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)

        assert '<aside class="speaker-notes"' not in html_out

    def test_full_presentation_pipeline(self):
        """Realistic multi-slide PPTX with all content types passes WCAG."""
        prs = Presentation()
        prs.core_properties.title = "NCDIT Report"

        # Slide 1: Title
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "Department Overview"

        # Slide 2: Content with body text
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Key Findings"
        s2.placeholders[1].text_frame.text = "Revenue increased by 15%."

        # Slide 3: Table
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "Budget"
        tbl_shape = s3.shapes.add_table(
            rows=3, cols=2,
            left=Inches(1), top=Inches(2),
            width=Inches(6), height=Inches(2),
        )
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Category"
        tbl.cell(0, 1).text = "Amount"
        tbl.cell(1, 0).text = "Personnel"
        tbl.cell(1, 1).text = "$1.5M"
        tbl.cell(2, 0).text = "Operations"
        tbl.cell(2, 1).text = "$0.5M"

        # Slide 4: Image
        s4 = prs.slides.add_slide(prs.slide_layouts[5])
        s4.shapes.title.text = "Architecture"
        s4.shapes.add_picture(
            io.BytesIO(_create_small_png()),
            left=Inches(2), top=Inches(2),
            width=Inches(4), height=Inches(3),
        )

        # Slide 5: Content with speaker notes
        s5 = prs.slides.add_slide(prs.slide_layouts[0])
        s5.shapes.title.text = "Next Steps"
        n5 = s5.notes_slide.notes_text_frame
        n5.text = "Wrap up and schedule follow-ups."

        pages, meta = extract_pptx(_pptx_bytes(prs))
        html_out, _ = build_html(pages, {}, meta)
        violations = validate_html(html_out)

        critical_serious = [
            v for v in violations if v.severity in ("critical", "serious")
        ]
        assert critical_serious == [], (
            f"Full pipeline violations: "
            f"{[(v.rule_id, v.description) for v in critical_serious]}"
        )

        # Spot-check content
        assert "Department Overview" in html_out
        assert "<table" in html_out
        assert "<figure>" in html_out
        assert "Speaker Notes" in html_out
        assert len(re.findall(r'<section\b', html_out)) == 5

    def test_sample_pptx_fixture_roundtrip(self):
        """If sample.pptx exists, verify the full roundtrip."""
        import os
        fixture = os.path.join(
            os.path.dirname(__file__), "..", "fixtures", "sample.pptx"
        )
        if not os.path.exists(fixture):
            pytest.skip("sample.pptx fixture not generated yet")

        with open(fixture, "rb") as f:
            data = f.read()

        pages, meta = extract_pptx(data)
        html_out, _ = build_html(pages, {}, meta)
        violations = validate_html(html_out)

        critical_serious = [
            v for v in violations if v.severity in ("critical", "serious")
        ]
        assert critical_serious == [], (
            f"sample.pptx roundtrip violations: "
            f"{[(v.rule_id, v.description) for v in critical_serious]}"
        )

        # Verify 5 slides → 5 sections
        sections = re.findall(r'<section\b', html_out)
        assert len(sections) == 5

        # Verify slide labels use "Slide" not "Page"
        assert "Slide 1" in html_out
        assert "Page 1" not in html_out
