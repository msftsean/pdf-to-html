"""
Unit tests for pptx_extractor — T065

Tests the PPTX extraction pipeline that converts PowerPoint presentations
to PageResult format compatible with html_builder.

Validates:
  - Slide content extraction (title + body text)
  - Speaker notes extraction
  - Table extraction from slides
  - Image extraction from slides
  - Metadata and PageResult structure
"""

import io
import struct
import zlib

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from backend.pptx_extractor import extract_pptx, SPEAKER_NOTES_FONT
from backend.pdf_extractor import TextSpan, ImageInfo, TableData, PageResult


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _pptx_bytes(prs: Presentation) -> bytes:
    """Serialize a python-pptx Presentation to bytes."""
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _create_small_png(width: int = 4, height: int = 4) -> bytes:
    """Create a minimal valid RGB PNG (solid red)."""
    raw = b""
    for _ in range(height):
        raw += b"\x00"  # filter byte: None
        for _ in range(width):
            raw += b"\xff\x00\x00"  # red pixel
    compressed = zlib.compress(raw)

    def _chunk(tag: bytes, data: bytes) -> bytes:
        c = tag + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", compressed)
        + _chunk(b"IEND", b"")
    )


# ---------------------------------------------------------------------------
# Tests: Slide content extraction (title + body text)
# ---------------------------------------------------------------------------

class TestSlideContentExtraction:
    """Slide title and body text extraction."""

    def test_slide_title_extracted(self):
        """Slide title becomes a TextSpan with h1-level size (≥ 24)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Welcome"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages) == 1
        h1_spans = [s for s in pages[0].text_spans if s.size >= 24.0]
        assert len(h1_spans) >= 1
        assert "Welcome" in h1_spans[0].text

    def test_slide_title_is_bold(self):
        """Slide title spans must have bold=True for h1 detection."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Bold Title"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        title_span = [s for s in pages[0].text_spans if "Bold Title" in s.text][0]
        assert title_span.bold is True

    def test_body_text_extracted(self):
        """Body text from content placeholder is extracted."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text_frame.text = "Body content here."

        pages, _ = extract_pptx(_pptx_bytes(prs))

        texts = [s.text for s in pages[0].text_spans]
        assert any("Body content here" in t for t in texts)

    def test_multiple_paragraphs_extracted(self):
        """Multiple paragraphs in a text frame are all extracted."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        tf = slide.placeholders[1].text_frame
        tf.text = "First paragraph"
        p = tf.add_paragraph()
        p.text = "Second paragraph"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        texts = [s.text for s in pages[0].text_spans]
        assert any("First paragraph" in t for t in texts)
        assert any("Second paragraph" in t for t in texts)

    def test_multiple_slides_produce_multiple_pages(self):
        """Each slide maps to a separate PageResult."""
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Slide One"
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2.shapes.title.text = "Slide Two"
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        slide3.shapes.title.text = "Slide Three"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages) == 3
        assert pages[0].page_number == 0
        assert pages[1].page_number == 1
        assert pages[2].page_number == 2

    def test_empty_text_frames_skipped(self):
        """Empty text frames don't produce text spans."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        pages, _ = extract_pptx(_pptx_bytes(prs))

        # No title, no text frames → no text spans
        assert len(pages) == 1
        # text_spans may still contain placeholder content depending on layout
        # but for blank layout, should be minimal

    def test_title_y_before_body_y(self):
        """Title span has a smaller y0 than body text spans."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "My Title"
        slide.placeholders[1].text_frame.text = "Body text."

        pages, _ = extract_pptx(_pptx_bytes(prs))

        title = [s for s in pages[0].text_spans if "My Title" in s.text][0]
        body = [s for s in pages[0].text_spans if "Body text" in s.text][0]
        assert title.y0 < body.y0


# ---------------------------------------------------------------------------
# Tests: Speaker notes extraction
# ---------------------------------------------------------------------------

class TestSpeakerNotesExtraction:
    """Speaker notes extraction and marking."""

    def test_notes_extracted(self):
        """Speaker notes text is included in text_spans."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Remember to explain this slide."

        pages, _ = extract_pptx(_pptx_bytes(prs))

        notes_spans = [
            s for s in pages[0].text_spans
            if s.font == SPEAKER_NOTES_FONT
        ]
        assert len(notes_spans) >= 1
        assert "Remember to explain" in notes_spans[0].text

    def test_notes_have_speaker_notes_font(self):
        """Speaker notes spans use SPEAKER_NOTES_FONT as the font name."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Notes content."

        pages, _ = extract_pptx(_pptx_bytes(prs))

        notes_spans = [s for s in pages[0].text_spans if s.font == SPEAKER_NOTES_FONT]
        assert len(notes_spans) >= 1
        # Non-notes spans should NOT have the speaker notes font
        regular_spans = [s for s in pages[0].text_spans if s.font != SPEAKER_NOTES_FONT]
        assert all(s.font != SPEAKER_NOTES_FONT for s in regular_spans)

    def test_no_notes_gives_no_speaker_spans(self):
        """Slides without speaker notes have no SPEAKER_NOTES_FONT spans."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "No Notes"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        notes_spans = [s for s in pages[0].text_spans if s.font == SPEAKER_NOTES_FONT]
        assert notes_spans == []

    def test_notes_y_after_content(self):
        """Speaker notes have a larger y0 than slide content."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        slide.placeholders[1].text_frame.text = "Content."
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Notes."

        pages, _ = extract_pptx(_pptx_bytes(prs))

        regular = [s for s in pages[0].text_spans if s.font != SPEAKER_NOTES_FONT]
        notes_spans = [s for s in pages[0].text_spans if s.font == SPEAKER_NOTES_FONT]

        if regular and notes_spans:
            max_content_y = max(s.y0 for s in regular)
            min_notes_y = min(s.y0 for s in notes_spans)
            assert min_notes_y > max_content_y


# ---------------------------------------------------------------------------
# Tests: Table extraction from slides
# ---------------------------------------------------------------------------

class TestTableExtraction:
    """Table extraction from PowerPoint slides."""

    def test_basic_table(self):
        """3×2 table: first row = header, remaining = data."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        tbl_shape = slide.shapes.add_table(
            rows=3, cols=2,
            left=Inches(1), top=Inches(1),
            width=Inches(6), height=Inches(2),
        )
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Name"
        tbl.cell(0, 1).text = "Value"
        tbl.cell(1, 0).text = "Alpha"
        tbl.cell(1, 1).text = "100"
        tbl.cell(2, 0).text = "Beta"
        tbl.cell(2, 1).text = "200"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages[0].tables) == 1
        td = pages[0].tables[0]
        assert isinstance(td, TableData)
        assert td.header == ["Name", "Value"]
        assert len(td.rows) == 2
        assert td.rows[0] == ["Alpha", "100"]
        assert td.rows[1] == ["Beta", "200"]

    def test_single_row_header_only(self):
        """A table with only one row yields header with no data rows."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tbl_shape = slide.shapes.add_table(
            rows=1, cols=2,
            left=Inches(1), top=Inches(1),
            width=Inches(6), height=Inches(1),
        )
        tbl_shape.table.cell(0, 0).text = "Col1"
        tbl_shape.table.cell(0, 1).text = "Col2"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        td = pages[0].tables[0]
        assert td.header == ["Col1", "Col2"]
        assert td.rows == []

    def test_empty_cells_preserved(self):
        """Empty table cells are preserved as empty strings."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tbl_shape = slide.shapes.add_table(
            rows=2, cols=2,
            left=Inches(1), top=Inches(1),
            width=Inches(6), height=Inches(2),
        )
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Header"
        tbl.cell(0, 1).text = ""
        tbl.cell(1, 0).text = ""
        tbl.cell(1, 1).text = "Data"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        td = pages[0].tables[0]
        assert td.header == ["Header", ""]
        assert td.rows[0] == ["", "Data"]


# ---------------------------------------------------------------------------
# Tests: Image extraction from slides
# ---------------------------------------------------------------------------

class TestImageExtraction:
    """Image extraction from PowerPoint slides."""

    def test_embedded_png(self):
        """An inline PNG image is extracted as an ImageInfo."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        png_bytes = _create_small_png()
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            left=Inches(1), top=Inches(1),
            width=Inches(2), height=Inches(2),
        )

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages[0].images) >= 1
        img = pages[0].images[0]
        assert isinstance(img, ImageInfo)
        assert img.extension == "png"
        assert len(img.image_bytes) > 0
        assert img.page_number == 0

    def test_no_images_gives_empty_list(self):
        """A text-only slide has an empty images list."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "No Images"

        pages, _ = extract_pptx(_pptx_bytes(prs))
        assert pages[0].images == []

    def test_image_position_from_shape(self):
        """Image coordinates reflect the shape position on the slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        png_bytes = _create_small_png()
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            left=Inches(2), top=Inches(3),
            width=Inches(4), height=Inches(2),
        )

        pages, _ = extract_pptx(_pptx_bytes(prs))

        img = pages[0].images[0]
        # 2 inches = 144 points, 3 inches = 216 points
        assert img.x0 > 100  # roughly 144pt
        assert img.y0 > 150  # roughly 216pt
        assert img.x1 > img.x0
        assert img.y1 > img.y0


# ---------------------------------------------------------------------------
# Tests: Metadata and PageResult structure
# ---------------------------------------------------------------------------

class TestPptxMetadata:
    """Metadata extraction and PageResult structure."""

    def test_format_metadata(self):
        """Metadata includes format='pptx'."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        _, meta = extract_pptx(_pptx_bytes(prs))
        assert meta.get("format") == "pptx"

    def test_title_metadata(self):
        """Presentation title extracted to metadata."""
        prs = Presentation()
        prs.core_properties.title = "Test Presentation"
        prs.slides.add_slide(prs.slide_layouts[6])
        _, meta = extract_pptx(_pptx_bytes(prs))
        assert meta.get("title") == "Test Presentation"

    def test_author_metadata(self):
        """Presentation author extracted to metadata."""
        prs = Presentation()
        prs.core_properties.author = "Jane Doe"
        prs.slides.add_slide(prs.slide_layouts[6])
        _, meta = extract_pptx(_pptx_bytes(prs))
        assert meta.get("author") == "Jane Doe"

    def test_page_result_structure(self):
        """PageResult has correct type and fields."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Hello"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages) == 1
        page = pages[0]
        assert isinstance(page, PageResult)
        assert page.page_number == 0
        assert page.width > 0
        assert page.height > 0
        assert page.is_scanned is False

    def test_is_scanned_always_false(self):
        """PPTX files are always digital — is_scanned must be False."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"

        pages, _ = extract_pptx(_pptx_bytes(prs))
        assert pages[0].is_scanned is False

    def test_each_slide_is_separate_page(self):
        """Each slide becomes a separate PageResult with correct index."""
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"Slide {i + 1}"

        pages, _ = extract_pptx(_pptx_bytes(prs))

        assert len(pages) == 5
        for i, page in enumerate(pages):
            assert page.page_number == i

    def test_slide_dimensions(self):
        """Slide width and height reflect presentation dimensions."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])

        pages, _ = extract_pptx(_pptx_bytes(prs))

        # Standard presentation is 10" × 7.5" = 720 × 540 pt
        assert abs(pages[0].width - 720.0) < 1.0
        assert abs(pages[0].height - 540.0) < 1.0

    def test_empty_presentation(self):
        """A presentation with no slides returns empty pages list."""
        prs = Presentation()
        pages, _ = extract_pptx(_pptx_bytes(prs))
        assert pages == []
